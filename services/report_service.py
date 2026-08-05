from __future__ import annotations

import json
import os
from dataclasses import dataclass
from datetime import date, datetime, timezone
from io import BytesIO
from typing import Any

import streamlit as st

from services.case_service import actor_name
from services.cyber_db import fetch_all, get_db, insert_row, update_rows
from services.trend_service import build_weekly_snapshot, fetch_news_for_analysis
from services.v6_audit_service import record_audit


@dataclass
class ReportNarrative:
    executive_summary: str
    trend_analysis: str
    priority_analysis: str
    recommendations: list[str]
    limitations: str
    source: str


def _format_period(snapshot: dict[str, Any]) -> str:
    return f"{snapshot['period']['start']} sampai {snapshot['period']['end']}"


def build_local_narrative(snapshot: dict[str, Any]) -> ReportNarrative:
    metrics = snapshot.get("metrics", {})
    top = snapshot.get("top_upt", {})
    total = int(metrics.get("total_publications") or 0)
    negative = int(metrics.get("negative_publications") or 0)
    mapped_negative = int(metrics.get("mapped_negative_publications") or 0)
    unmapped_negative = int(metrics.get("unmapped_negative_publications") or 0)
    upt_count = int(metrics.get("negative_upt_count") or 0)
    media = int(metrics.get("unique_media") or 0)
    issues = int(metrics.get("issue_count") or 0)
    high = int(metrics.get("high_critical_count") or 0)
    change = metrics.get("negative_change_percent")
    concentration = float(metrics.get("top_two_concentration_percent") or 0)
    top_name = str(top.get("UPT") or "Belum tersedia")
    top_count = int(top.get("Berita Negatif") or top.get("Jumlah Publikasi") or 0)
    top_issue = str(top.get("Isu Utama") or "Belum dapat ditentukan")

    executive = (
        f"Dalam periode {_format_period(snapshot)}, sistem menghimpun {total} publikasi unik. "
        f"Sebanyak {negative} publikasi diklasifikasikan negatif. Dari jumlah tersebut, {mapped_negative} publikasi dapat dipetakan kepada {upt_count} UPT"
        f" dan {unmapped_negative} publikasi masih menunggu pemetaan UPT. "
        f"Pemberitaan berasal dari {media} media dan membentuk sekitar {issues} kelompok isu. "
        f"UPT dengan sorotan negatif tertinggi adalah {top_name} dengan {top_count} publikasi."
    )
    if change is None:
        trend = (
            "Minggu sebelumnya tidak memiliki publikasi negatif yang sebanding, sehingga persentase perubahan tidak dihitung. "
            "Kenaikan perlu dibaca sebagai kemunculan eksposur baru, bukan otomatis sebagai pertambahan jumlah kejadian."
        )
    elif change > 0:
        trend = (
            f"Pemberitaan negatif meningkat {abs(change):.1f} persen dibanding periode sebelumnya. "
            f"Dua UPT teratas menyumbang {concentration:.1f} persen dari seluruh publikasi negatif, "
            "sehingga konsentrasi eksposur perlu dibedakan dari jumlah kejadian faktual."
        )
    elif change < 0:
        trend = (
            f"Pemberitaan negatif menurun {abs(change):.1f} persen dibanding periode sebelumnya. "
            "Penurunan kuantitas tidak meniadakan kebutuhan tindak lanjut pada isu berurgensi tinggi atau kritis."
        )
    else:
        trend = "Jumlah publikasi negatif relatif tetap dibanding periode sebelumnya. Pergeseran media, UPT, dan urgensi tetap perlu diperiksa."

    priority = (
        f"Isu utama pada UPT paling disorot adalah: {top_issue}. "
        f"Terdapat {high} publikasi berurgensi tinggi atau kritis. Publikasi berulang mengenai satu kasus tetap dihitung sebagai eksposur media, bukan sebagai kejadian baru yang berbeda."
    )
    recommendations = [
        "Dahulukan telaah dan klarifikasi pada UPT dengan publikasi negatif serta urgensi tertinggi.",
        "Pisahkan kejadian baru, perkembangan kasus lama, dan konten lama yang kembali viral sebelum menyusun bahan pimpinan.",
        "Pastikan setiap angka pada laporan dapat ditelusuri kembali ke daftar link unik yang menjadi sumbernya.",
    ]
    if high:
        recommendations.insert(0, "Tetapkan penanggung jawab dan tenggat tindak lanjut untuk seluruh berita tinggi atau kritis.")
    limitations = (
        "Laporan menghitung publikasi berdasarkan link unik. Kesamaan isu tidak dianggap sebagai duplikasi. "
        "Pengelompokan isu otomatis merupakan bantuan awal dan tetap memerlukan validasi Analis Intelijen Pemberitaan."
    )
    return ReportNarrative(executive, trend, priority, recommendations, limitations, "local_rules")


def _extract_openai_text(response: Any) -> str:
    output_text = getattr(response, "output_text", None)
    if output_text:
        return str(output_text)
    choices = getattr(response, "choices", None)
    if choices:
        message = getattr(choices[0], "message", None)
        content = getattr(message, "content", None)
        if content:
            return str(content)
    return ""


def build_ai_narrative(snapshot: dict[str, Any]) -> ReportNarrative:
    try:
        secret_api_key = st.secrets.get("OPENAI_API_KEY", "")
        secret_model = st.secrets.get("OPENAI_MODEL", "")
    except Exception:
        secret_api_key = ""
        secret_model = ""
    api_key = str(os.getenv("OPENAI_API_KEY") or secret_api_key or "").strip()
    if not api_key:
        return build_local_narrative(snapshot)
    try:
        from openai import OpenAI

        client = OpenAI(api_key=api_key)
        model = str(os.getenv("OPENAI_MODEL") or secret_model or "gpt-5-mini").strip() or "gpt-5-mini"
        compact = {
            "period": snapshot.get("period"),
            "previous_period": snapshot.get("previous_period"),
            "metrics": snapshot.get("metrics"),
            "top_upt": snapshot.get("top_upt"),
            "upt_table": snapshot.get("upt_table", [])[:15],
            "urgency_distribution": snapshot.get("urgency_distribution", []),
            "daily_trend": snapshot.get("daily_trend", []),
        }
        system_instruction = (
            "Anda menyusun laporan intelijen pemberitaan Pemasyarakatan. Gunakan hanya angka dalam payload. "
            "Jangan mengubah jumlah publikasi menjadi jumlah kejadian. Link unik adalah satu publikasi, sedangkan isu yang sama boleh memiliki banyak publikasi. "
            "Keluaran wajib JSON valid dengan kunci executive_summary, trend_analysis, priority_analysis, recommendations, limitations. "
            "recommendations harus berupa array 3 sampai 5 butir. Gunakan Bahasa Indonesia formal, padat, dan tidak membuat tuduhan baru."
        )
        user_prompt = "Susun narasi laporan mingguan dari payload berikut:\n" + json.dumps(compact, ensure_ascii=False)

        try:
            response = client.responses.create(
                model=model,
                instructions=system_instruction,
                input=user_prompt,
            )
        except Exception:
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.2,
            )
        text = _extract_openai_text(response).strip()
        if text.startswith("```"):
            text = text.strip("`")
            if text.lower().startswith("json"):
                text = text[4:].strip()
        parsed = json.loads(text)
        return ReportNarrative(
            str(parsed.get("executive_summary") or ""),
            str(parsed.get("trend_analysis") or ""),
            str(parsed.get("priority_analysis") or ""),
            [str(item) for item in parsed.get("recommendations") or []],
            str(parsed.get("limitations") or ""),
            "openai",
        )
    except Exception:
        fallback = build_local_narrative(snapshot)
        return ReportNarrative(
            fallback.executive_summary,
            fallback.trend_analysis,
            fallback.priority_analysis,
            fallback.recommendations,
            fallback.limitations,
            "local_fallback_after_ai_error",
        )


def build_report_package(start: date, end: date, use_ai: bool = True) -> tuple[dict[str, Any], ReportNarrative]:
    rows = fetch_news_for_analysis()
    snapshot = build_weekly_snapshot(rows, start, end)
    narrative = build_ai_narrative(snapshot) if use_ai else build_local_narrative(snapshot)
    return snapshot, narrative


def save_weekly_report(
    snapshot: dict[str, Any],
    narrative: ReportNarrative,
    user: Any,
    status: str = "Draf Sistem",
) -> dict[str, Any]:
    payload = {
        "period_start": snapshot["period"]["start"],
        "period_end": snapshot["period"]["end"],
        "status": status,
        "snapshot_data": snapshot,
        "ai_narrative": {
            "executive_summary": narrative.executive_summary,
            "trend_analysis": narrative.trend_analysis,
            "priority_analysis": narrative.priority_analysis,
            "recommendations": narrative.recommendations,
            "limitations": narrative.limitations,
            "source": narrative.source,
        },
        "created_by": actor_name(user),
        "ai_provider": narrative.source,
    }
    row = insert_row("weekly_reports", payload)
    record_audit(user, "weekly_report.create", "weekly_report", str(row.get("id") or ""), {"period_start": snapshot["period"]["start"], "period_end": snapshot["period"]["end"], "status": status, "ai_provider": narrative.source})
    return row


def fetch_weekly_reports(limit: int = 200) -> list[dict[str, Any]]:
    return fetch_all("weekly_reports", "*", order_by="created_at", desc=True, max_rows=limit)


def update_weekly_report_status(report_id: str, status: str, user: Any) -> None:
    payload: dict[str, Any] = {"status": status, "updated_by": actor_name(user)}
    now = datetime.now(timezone.utc).isoformat()
    if status == "Diverifikasi":
        payload.update({"verified_by": actor_name(user), "verified_at": now})
    if status == "Disetujui":
        payload.update({"approved_by": actor_name(user), "approved_at": now})
    if status == "Dipublikasikan":
        payload.update({"published_by": actor_name(user), "published_at": now, "locked_at": now})
    update_rows("weekly_reports", payload, filters=[("eq", "id", report_id)])
    record_audit(user, "weekly_report.status_update", "weekly_report", report_id, {"status": status})


def _title(snapshot: dict[str, Any]) -> str:
    return "Laporan Intelijen Pemberitaan Mingguan"


def generate_pdf(snapshot: dict[str, Any], narrative: ReportNarrative) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    buffer = BytesIO()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="CenterTitle", parent=styles["Title"], alignment=TA_CENTER, spaceAfter=18))
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=1.5 * cm, leftMargin=1.5 * cm, topMargin=1.4 * cm, bottomMargin=1.4 * cm)
    story: list[Any] = []
    story.append(Paragraph(_title(snapshot), styles["CenterTitle"]))
    story.append(Paragraph(f"Periode {_format_period(snapshot)}", styles["Heading2"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Ringkasan Eksekutif", styles["Heading1"]))
    story.append(Paragraph(narrative.executive_summary, styles["BodyText"]))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Analisis Tren", styles["Heading2"]))
    story.append(Paragraph(narrative.trend_analysis, styles["BodyText"]))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Analisis Prioritas", styles["Heading2"]))
    story.append(Paragraph(narrative.priority_analysis, styles["BodyText"]))

    metrics = snapshot.get("metrics", {})
    metric_data = [
        ["Indikator", "Nilai"],
        ["Total publikasi unik", metrics.get("total_publications", 0)],
        ["Publikasi negatif", metrics.get("negative_publications", 0)],
        ["Negatif terpetakan", metrics.get("mapped_negative_publications", 0)],
        ["Negatif belum terpetakan", metrics.get("unmapped_negative_publications", 0)],
        ["UPT terpetakan", metrics.get("negative_upt_count", 0)],
        ["Media unik", metrics.get("unique_media", 0)],
        ["Kelompok isu", metrics.get("issue_count", 0)],
        ["Urgensi tinggi/kritis", metrics.get("high_critical_count", 0)],
    ]
    story.append(Spacer(1, 14))
    metric_table = Table(metric_data, colWidths=[10 * cm, 5 * cm])
    metric_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8E8E8")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(metric_table)

    story.append(PageBreak())
    story.append(Paragraph("Rekap Publikasi Negatif per UPT", styles["Heading1"]))
    upt_rows = snapshot.get("upt_table", [])
    table_data = [["No", "UPT", "Publikasi", "Media", "Isu", "Negatif", "Urgensi"]]
    for idx, row in enumerate(upt_rows[:30], 1):
        table_data.append([
            idx,
            Paragraph(str(row.get("UPT") or ""), styles["BodyText"]),
            row.get("Jumlah Publikasi", 0),
            row.get("Jumlah Media", 0),
            row.get("Jumlah Isu", 0),
            row.get("Berita Negatif", 0),
            row.get("Urgensi Tertinggi", "Rendah"),
        ])
    upt_table = Table(table_data, repeatRows=1, colWidths=[0.7 * cm, 8.2 * cm, 1.5 * cm, 1.2 * cm, 1.1 * cm, 1.2 * cm, 1.8 * cm])
    upt_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8E8E8")),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(upt_table)

    story.append(Spacer(1, 14))
    story.append(Paragraph("Rekomendasi", styles["Heading1"]))
    for idx, item in enumerate(narrative.recommendations, 1):
        story.append(Paragraph(f"{idx}. {item}", styles["BodyText"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Catatan Batasan", styles["Heading2"]))
    story.append(Paragraph(narrative.limitations, styles["BodyText"]))
    doc.build(story)
    return buffer.getvalue()


def generate_docx(snapshot: dict[str, Any], narrative: ReportNarrative) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    buffer = BytesIO()
    document = Document()
    title = document.add_heading(_title(snapshot), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    period = document.add_paragraph(f"Periode {_format_period(snapshot)}")
    period.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_heading("Ringkasan Eksekutif", level=1)
    document.add_paragraph(narrative.executive_summary)
    document.add_heading("Analisis Tren", level=1)
    document.add_paragraph(narrative.trend_analysis)
    document.add_heading("Analisis Prioritas", level=1)
    document.add_paragraph(narrative.priority_analysis)

    document.add_heading("Indikator Utama", level=1)
    metrics = snapshot.get("metrics", {})
    metric_table = document.add_table(rows=1, cols=2)
    metric_table.style = "Table Grid"
    metric_table.rows[0].cells[0].text = "Indikator"
    metric_table.rows[0].cells[1].text = "Nilai"
    for label, value in [
        ("Total publikasi unik", metrics.get("total_publications", 0)),
        ("Publikasi negatif", metrics.get("negative_publications", 0)),
        ("Negatif terpetakan", metrics.get("mapped_negative_publications", 0)),
        ("Negatif belum terpetakan", metrics.get("unmapped_negative_publications", 0)),
        ("UPT terpetakan", metrics.get("negative_upt_count", 0)),
        ("Media unik", metrics.get("unique_media", 0)),
        ("Kelompok isu", metrics.get("issue_count", 0)),
        ("Urgensi tinggi/kritis", metrics.get("high_critical_count", 0)),
    ]:
        cells = metric_table.add_row().cells
        cells[0].text = label
        cells[1].text = str(value)

    document.add_heading("Rekap Publikasi Negatif per UPT", level=1)
    upt_table = document.add_table(rows=1, cols=7)
    upt_table.style = "Table Grid"
    headers = ["No", "UPT", "Publikasi", "Media", "Isu", "Negatif", "Urgensi"]
    for cell, text in zip(upt_table.rows[0].cells, headers):
        cell.text = text
    for idx, row in enumerate(snapshot.get("upt_table", [])[:50], 1):
        cells = upt_table.add_row().cells
        values = [
            idx, row.get("UPT", ""), row.get("Jumlah Publikasi", 0), row.get("Jumlah Media", 0),
            row.get("Jumlah Isu", 0), row.get("Berita Negatif", 0), row.get("Urgensi Tertinggi", "Rendah"),
        ]
        for cell, value in zip(cells, values):
            cell.text = str(value)

    document.add_heading("Rekomendasi", level=1)
    for item in narrative.recommendations:
        document.add_paragraph(item, style="List Number")
    document.add_heading("Catatan Batasan", level=1)
    document.add_paragraph(narrative.limitations)
    document.add_heading("Lampiran Publikasi Prioritas", level=1)
    for item in snapshot.get("top_news", [])[:20]:
        paragraph = document.add_paragraph(style="List Bullet")
        run = paragraph.add_run(str(item.get("judul") or "Tanpa judul"))
        run.bold = True
        paragraph.add_run(f"\nUPT: {item.get('nama_upt') or 'Belum Teridentifikasi'} | Media: {item.get('media') or 'Tidak diketahui'} | Urgensi: {item.get('urgensi') or 'Rendah'}")
        if item.get("link_normalized"):
            paragraph.add_run(f"\nLink: {item['link_normalized']}")

    for section in document.sections:
        section.header.paragraphs[0].text = "CYBER-INTELPAS"
        section.footer.paragraphs[0].text = "Dokumen otomatis, wajib diverifikasi sebelum digunakan sebagai laporan resmi."
    document.save(buffer)
    return buffer.getvalue()


def generate_pptx(snapshot: dict[str, Any], narrative: ReportNarrative) -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt

    buffer = BytesIO()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, title: str, subtitle: str = ""):
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(slide, _title(snapshot), f"Periode {_format_period(snapshot)}\nCYBER-INTELPAS")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(slide, "Executive Summary")
    slide.placeholders[1].text = narrative.executive_summary

    metrics = snapshot.get("metrics", {})
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "KPI Pemberitaan Mingguan"
    metric_items = [
        ("Total Publikasi", metrics.get("total_publications", 0)),
        ("Negatif Terpetakan", metrics.get("mapped_negative_publications", 0)),
        ("UPT", metrics.get("negative_upt_count", 0)),
        ("Media", metrics.get("unique_media", 0)),
        ("Isu", metrics.get("issue_count", 0)),
        ("Tinggi/Kritis", metrics.get("high_critical_count", 0)),
    ]
    for idx, (label, value) in enumerate(metric_items):
        col, row = idx % 3, idx // 3
        box = slide.shapes.add_textbox(Inches(0.7 + col * 4.2), Inches(1.5 + row * 2.4), Inches(3.5), Inches(1.6))
        frame = box.text_frame
        p = frame.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(30)
        p.font.bold = True
        p2 = frame.add_paragraph()
        p2.text = label
        p2.font.size = Pt(15)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Tren Publikasi Harian"
    chart_data = CategoryChartData()
    trend = snapshot.get("daily_trend", [])
    chart_data.categories = [str(item.get("Tanggal")) for item in trend]
    chart_data.add_series("Publikasi", [int(item.get("Jumlah Publikasi") or 0) for item in trend])
    slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.9), chart_data)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "UPT dengan Eksposur Negatif Tertinggi"
    top_rows = snapshot.get("upt_table", [])[:10]
    chart_data = CategoryChartData()
    chart_data.categories = [str(item.get("UPT") or "")[:35] for item in top_rows]
    chart_data.add_series("Publikasi Negatif", [int(item.get("Berita Negatif") or 0) for item in top_rows])
    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.2), chart_data)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(slide, "Analisis Tren")
    slide.placeholders[1].text = narrative.trend_analysis

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(slide, "Isu Prioritas")
    slide.placeholders[1].text = narrative.priority_analysis

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(slide, "Rekomendasi")
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for idx, item in enumerate(narrative.recommendations):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title(slide, "Catatan Validasi")
    slide.placeholders[1].text = narrative.limitations + "\n\nKeputusan akhir dan pengesahan tetap berada pada pejabat yang berwenang."

    prs.save(buffer)
    return buffer.getvalue()
